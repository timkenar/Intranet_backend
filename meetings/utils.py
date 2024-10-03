import os
from fpdf import FPDF
from PyPDF2 import PdfMerger
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from odf.opendocument import OpenDocumentSpreadsheet
from odf.table import Table
from odf.text import P
from django.conf import settings

def convert_docx_to_pdf(input_file, output_file):
    """Convert .docx to PDF."""
    doc = DocxDocument(input_file)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font('Arial', 'B', 12)
    pdf.cell(0, 10, 'Document: ' + os.path.basename(input_file), ln=True)

    pdf.set_font('Arial', '', 12)
    for para in doc.paragraphs:
        pdf.multi_cell(0, 10, para.text)

    pdf.output(output_file)
    return output_file

def convert_pptx_to_pdf(input_file, output_file):
    """Convert .pptx to PDF."""
    presentation = Presentation(input_file)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    for slide in presentation.slides:
        pdf.add_page()
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                pdf.set_font('Arial', 'B', 12)
                pdf.multi_cell(0, 10, shape.text)

    pdf.output(output_file)
    return output_file

def convert_xlsx_to_pdf(input_file, output_file):
    """Convert .xlsx to PDF."""
    workbook = load_workbook(input_file)
    sheet = workbook.active
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font('Arial', 'B', 12)
    pdf.cell(0, 10, 'Excel Sheet: ' + os.path.basename(input_file), ln=True)

    pdf.set_font('Arial', '', 12)
    for row in sheet.iter_rows(values_only=True):
        row_data = ' | '.join([str(cell) if cell is not None else '' for cell in row])
        pdf.multi_cell(0, 10, row_data)

    pdf.output(output_file)
    return output_file

def convert_ods_to_pdf(input_file, output_file):
    """Convert .ods to PDF."""
    odoc = OpenDocumentSpreadsheet(input_file)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font('Arial', 'B', 12)
    pdf.cell(0, 10, 'Spreadsheet: ' + os.path.basename(input_file), ln=True)

    pdf.set_font('Arial', '', 12)
    for table in odoc.getElementsByType(Table):
        for row in table.getElementsByType(P):
            row_data = ' | '.join([cell.firstChild.data for cell in row.childNodes if cell.firstChild])
            pdf.multi_cell(0, 10, row_data)

    pdf.output(output_file)
    return output_file

def merge_meeting_documents(meeting):
    """Merge all documents related to a meeting into a single boardpack PDF."""
    merger = PdfMerger()

    # Create a temp directory for storing converted PDFs
    temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp_pdfs')
    os.makedirs(temp_dir, exist_ok=True)

    agenda_items = meeting.agenda_items.all()

    for agenda_item in agenda_items:
        documents = agenda_item.documents.all()
        for document in documents:
            doc_path = document.file.path
            # Output PDF path for each converted document
            output_pdf = os.path.join(temp_dir, f"{document.pk}.pdf")

            # Handle different file types
            if doc_path.endswith('.pdf'):
                merger.append(doc_path)
            elif doc_path.endswith('.docx'):
                convert_docx_to_pdf(doc_path, output_pdf)
                merger.append(output_pdf)
            elif doc_path.endswith('.pptx'):
                convert_pptx_to_pdf(doc_path, output_pdf)
                merger.append(output_pdf)
            elif doc_path.endswith('.xlsx'):
                convert_xlsx_to_pdf(doc_path, output_pdf)
                merger.append(output_pdf)
            elif doc_path.endswith('.ods'):
                convert_ods_to_pdf(doc_path, output_pdf)
                merger.append(output_pdf)

    # Define the output path for the merged boardpack
    boardpack_filename = f"Boardpack_{meeting.title.replace(' ', '_')}.pdf"
    boardpack_path = os.path.join(settings.MEDIA_ROOT, 'boardpacks', boardpack_filename)

    # Write the merged PDF to the output file
    with open(boardpack_path, 'wb') as boardpack_pdf:
        merger.write(boardpack_pdf)

    # Close the PdfMerger
    merger.close()

    # Cleanup temporary PDFs
    for file in os.listdir(temp_dir):
        os.remove(os.path.join(temp_dir, file))

    return boardpack_path
